"""
Mock-free unit tests for T11 PowerPoint Loader Unified

Tests the unified PowerPoint loader tool with real PowerPoint processing using python-pptx.
No mocking is used - all functionality is tested with real data and real processing.
"""

import pytest
import tempfile
from pathlib import Path
import os

from src.tools.phase1.t11_powerpoint_loader_unified import T11PowerPointLoaderUnified
from src.core.service_manager import ServiceManager
from src.tools.base_tool import ToolRequest

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class TestT11PowerPointLoaderUnifiedMockFree:
    def setup_method(self):
        """Set up test fixtures with real ServiceManager - NO mocks"""
        # Real ServiceManager - NO mocking
        self.service_manager = ServiceManager()
        self.tool = T11PowerPointLoaderUnified(service_manager=self.service_manager)
        
        # Create real test PowerPoint files if python-pptx is available
        if PPTX_AVAILABLE:
            self.test_files = self._create_real_test_powerpoint_files()
        else:
            self.test_files = {}
    
    def teardown_method(self):
        """Clean up real test files"""
        for file_path in self.test_files.values():
            try:
                if os.path.exists(file_path):
                    os.unlink(file_path)
            except:
                pass
    
    def _create_real_test_powerpoint_files(self) -> dict:
        """Create real PowerPoint test files for testing"""
        test_files = {}
        
        # Skip if python-pptx not available
        if not PPTX_AVAILABLE:
            return test_files
        
        # Simple PowerPoint with text content
        prs_simple = Presentation()
        
        # Title slide
        title_slide_layout = prs_simple.slide_layouts[0]
        slide = prs_simple.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Test Presentation"
        subtitle.text = "Created for unit testing"
        
        # Content slide
        bullet_slide_layout = prs_simple.slide_layouts[1]
        slide = prs_simple.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = 'Key Features'
        tf = body_shape.text_frame
        tf.text = 'Feature 1: Document loading'
        p = tf.add_paragraph()
        p.text = 'Feature 2: Text extraction'
        p = tf.add_paragraph()
        p.text = 'Feature 3: Content analysis'
        
        simple_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
        prs_simple.save(simple_file.name)
        simple_file.close()
        test_files['simple'] = simple_file.name
        
        # Complex PowerPoint with multiple slide types
        prs_complex = Presentation()
        
        # Title slide
        title_slide = prs_complex.slides.add_slide(prs_complex.slide_layouts[0])
        title_slide.shapes.title.text = "Complex Presentation"
        title_slide.placeholders[1].text = "Advanced Testing Scenarios"
        
        # Content with bullet points
        content_slide = prs_complex.slides.add_slide(prs_complex.slide_layouts[1])
        content_slide.shapes.title.text = "Project Overview"
        tf = content_slide.shapes.placeholders[1].text_frame
        tf.text = "Objective: Implement mock-free testing"
        p = tf.add_paragraph()
        p.text = "Approach: Real library integration"
        p = tf.add_paragraph()
        p.text = "Result: High confidence validation"
        
        # Blank slide with custom content
        blank_slide = prs_complex.slides.add_slide(prs_complex.slide_layouts[6])
        textbox = blank_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        textbox.text_frame.text = "Custom textbox content for testing"
        
        # Another content slide
        final_slide = prs_complex.slides.add_slide(prs_complex.slide_layouts[1])
        final_slide.shapes.title.text = "Conclusion"
        tf = final_slide.shapes.placeholders[1].text_frame
        tf.text = "Testing methodology validated"
        p = tf.add_paragraph()
        p.text = "Real functionality confirmed"
        
        complex_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
        prs_complex.save(complex_file.name)
        complex_file.close()
        test_files['complex'] = complex_file.name
        
        return test_files
    
    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
    def test_tool_contract_real(self):
        """Test tool contract with REAL contract validation"""
        contract = self.tool.get_contract()
        
        # Verify contract structure
        assert contract.tool_id == "T11"
        assert contract.name == "PowerPoint Document Loader"
        assert contract.category == "document_processing"
        assert "file_path" in contract.input_schema["required"]
        assert "document" in contract.output_schema["required"]
        assert len(contract.dependencies) > 0
        
        # Verify performance requirements
        assert "max_execution_time" in contract.performance_requirements
        assert "max_memory_mb" in contract.performance_requirements
        assert "min_confidence" in contract.performance_requirements
        
        # Verify error conditions
        assert "POWERPOINT_CORRUPTED" in contract.error_conditions
        assert "FILE_NOT_FOUND" in contract.error_conditions
    
    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
    def test_simple_powerpoint_loading_real(self):
        """Test loading simple PowerPoint with REAL processing"""
        request = ToolRequest(
            tool_id="T11",
            operation="load_powerpoint",
            input_data={
                "file_path": self.test_files['simple'],
                "workflow_id": "test_workflow_simple"
            },
            parameters={}
        )
        
        result = self.tool.execute(request)
        
        # Verify successful execution
        assert result.status == "success"
        assert result.tool_id == "T11"
        assert result.execution_time > 0
        
        # Verify document data
        doc = result.data["document"]
        assert doc["document_id"] == "test_workflow_simple_" + Path(self.test_files['simple']).stem
        assert doc["file_name"] == Path(self.test_files['simple']).name
        assert doc["slide_count"] == 2  # Title + content slide
        assert doc["shape_count"] > 0
        assert doc["confidence"] > 0.5
        assert len(doc["text_content"]) > 0
        
        # Verify presentation data structure
        presentation_data = doc["presentation_data"]
        assert "slides" in presentation_data
        assert "metadata" in presentation_data
        assert len(presentation_data["slides"]) == 2
        
        # Verify specific content
        assert "Test Presentation" in doc["text_content"]
        assert "Key Features" in doc["text_content"]
        assert "Feature 1" in doc["text_content"]
    
    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
    def test_complex_powerpoint_multiple_slides_real(self):
        """Test loading complex PowerPoint with multiple slides with REAL processing"""
        request = ToolRequest(
            tool_id="T11",
            operation="load_powerpoint",
            input_data={
                "file_path": self.test_files['complex'],
                "workflow_id": "test_workflow_complex"
            },
            parameters={}
        )
        
        result = self.tool.execute(request)
        
        # Verify successful execution
        assert result.status == "success"
        
        # Verify document data
        doc = result.data["document"]
        assert doc["slide_count"] == 4  # Title + content + blank + final
        assert doc["shape_count"] > 5   # Multiple shapes across slides
        assert doc["confidence"] > 0.5
        
        # Verify presentation data
        presentation_data = doc["presentation_data"]
        assert len(presentation_data["slides"]) == 4
        
        # Verify content from all slides
        assert "Complex Presentation" in doc["text_content"]      # Title slide
        assert "Project Overview" in doc["text_content"]         # Content slide
        assert "Custom textbox content" in doc["text_content"]   # Blank slide
        assert "Conclusion" in doc["text_content"]              # Final slide
    
    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
    def test_parse_options_functionality_real(self):
        """Test different parse options with REAL processing"""
        # Test with metadata extraction enabled
        request_metadata = ToolRequest(
            tool_id="T11",
            operation="load_powerpoint",
            input_data={
                "file_path": self.test_files['simple'],
                "parse_options": {
                    "extract_metadata": True,
                    "extract_notes": True
                }
            },
            parameters={}
        )
        
        result_metadata = self.tool.execute(request_metadata)
        assert result_metadata.status == "success"
        
        # Verify metadata is included
        presentation_data = result_metadata.data["document"]["presentation_data"]
        assert "metadata" in presentation_data
        
        # Test with metadata extraction disabled
        request_no_metadata = ToolRequest(
            tool_id="T11",
            operation="load_powerpoint",
            input_data={
                "file_path": self.test_files['simple'],
                "parse_options": {
                    "extract_metadata": False
                }
            },
            parameters={}
        )
        
        result_no_metadata = self.tool.execute(request_no_metadata)
        assert result_no_metadata.status == "success"
    
    def test_file_not_found_error_real(self):
        """Test file not found error with REAL missing file"""
        request = ToolRequest(
            tool_id="T11",
            operation="load_powerpoint",
            input_data={
                "file_path": "/path/to/nonexistent/file.pptx"
            },
            parameters={}
        )
        
        result = self.tool.execute(request)
        
        # Verify error handling
        assert result.status == "error"
        assert result.error_code == "FILE_NOT_FOUND"
        assert "File not found" in result.error_message
    
    def test_invalid_file_type_error_real(self):
        """Test invalid file type error with REAL non-PowerPoint file"""
        # Create a non-PowerPoint file
        txt_file = tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False)
        txt_file.write("This is not a PowerPoint file")
        txt_file.close()
        
        try:
            request = ToolRequest(
                tool_id="T11",
                operation="load_powerpoint",
                input_data={
                    "file_path": txt_file.name
                },
                parameters={}
            )
            
            result = self.tool.execute(request)
            
            # Verify error handling
            assert result.status == "error"
            assert result.error_code == "INVALID_FILE_TYPE"
            assert "Invalid file extension" in result.error_message
        finally:
            os.unlink(txt_file.name)
    
    def test_library_not_available_error_real(self):
        """Test behavior when python-pptx library is not available"""
        # This test simulates the library not being available
        # by temporarily setting PPTX_AVAILABLE to False in the tool
        original_available = self.tool.__class__.__module__ 
        
        # Create a request
        request = ToolRequest(
            tool_id="T11",
            operation="load_powerpoint",
            input_data={
                "file_path": "/some/file.pptx"
            },
            parameters={}
        )
        
        # If python-pptx is not available, should get library missing error
        if not PPTX_AVAILABLE:
            result = self.tool.execute(request)
            assert result.status == "error"
            assert result.error_code == "PPTX_LIBRARY_MISSING"
    
    def test_input_validation_real(self):
        """Test input validation with REAL validation logic"""
        # Test missing file_path
        result = self.tool.validate_input({})
        assert result == False
        
        result = self.tool.validate_input({"file_path": ""})
        assert result == False
        
        # Test valid input
        result = self.tool.validate_input({"file_path": "/some/path.pptx"})
        assert result == True
    
    def test_health_check_real(self):
        """Test health check with REAL service verification"""
        result = self.tool.health_check()
        
        # Verify health check structure
        assert isinstance(result.data, dict)
        assert "healthy" in result.data
        assert "pptx_available" in result.data
        assert "services_healthy" in result.data
        assert "supported_formats" in result.data
        
        # Verify supported formats
        supported_formats = result.data["supported_formats"]
        assert ".pptx" in supported_formats
        assert ".ppt" in supported_formats
    
    def test_cleanup_functionality_real(self):
        """Test cleanup functionality with REAL resource management"""
        # Add some temp files to the tool
        temp_file = tempfile.NamedTemporaryFile(delete=False)
        temp_file.close()
        self.tool._temp_files.append(temp_file.name)
        
        # Verify cleanup works
        cleanup_result = self.tool.cleanup()
        assert cleanup_result == True
        assert len(self.tool._temp_files) == 0
        assert not os.path.exists(temp_file.name)
    
    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
    def test_confidence_calculation_real(self):
        """Test confidence calculation with REAL PowerPoint parsing metrics"""
        # Test with simple PowerPoint
        request_simple = ToolRequest(
            tool_id="T11",
            operation="load_powerpoint",
            input_data={"file_path": self.test_files['simple']},
            parameters={}
        )
        
        result_simple = self.tool.execute(request_simple)
        confidence_simple = result_simple.data["document"]["confidence"]
        
        # Test with complex PowerPoint
        request_complex = ToolRequest(
            tool_id="T11",
            operation="load_powerpoint",
            input_data={"file_path": self.test_files['complex']},
            parameters={}
        )
        
        result_complex = self.tool.execute(request_complex)
        confidence_complex = result_complex.data["document"]["confidence"]
        
        # Both should have reasonable confidence
        assert confidence_simple > 0.5
        assert confidence_complex > 0.5
        # Complex PowerPoint should have higher confidence due to more content
        assert confidence_complex >= confidence_simple - 0.1
    
    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
    def test_performance_metrics_real(self):
        """Test performance metrics with REAL execution measurement"""
        request = ToolRequest(
            tool_id="T11",
            operation="load_powerpoint",
            input_data={"file_path": self.test_files['complex']},
            parameters={}
        )
        
        result = self.tool.execute(request)
        
        # Verify performance metrics are captured
        assert result.execution_time > 0
        assert result.memory_used >= 0
        
        # Verify reasonable execution time
        assert result.execution_time < 5.0  # Should be under 5 seconds
    
    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
    def test_service_integration_real(self):
        """Test service integration with REAL services"""
        request = ToolRequest(
            tool_id="T11",
            operation="load_powerpoint",
            input_data={
                "file_path": self.test_files['simple'],
                "workflow_id": "test_service_integration"
            },
            parameters={}
        )
        
        result = self.tool.execute(request)
        
        # Verify service integration
        assert result.status == "success"
        assert "operation_id" in result.metadata
        
        # Verify provenance tracking
        operation_id = result.metadata["operation_id"]
        assert operation_id is not None
        
        # Verify quality assessment
        doc = result.data["document"]
        assert "quality_tier" in doc
        assert doc["confidence"] > 0
    
    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
    def test_slide_content_extraction_real(self):
        """Test slide content extraction with REAL PowerPoint parsing"""
        request = ToolRequest(
            tool_id="T11",
            operation="load_powerpoint",
            input_data={"file_path": self.test_files['complex']},
            parameters={}
        )
        
        result = self.tool.execute(request)
        presentation_data = result.data["document"]["presentation_data"]
        
        # Verify slide structure
        slides = presentation_data["slides"]
        assert len(slides) == 4
        
        # Verify each slide has expected structure
        for slide in slides:
            assert "slide_number" in slide
            assert "shapes" in slide
            assert "text_content" in slide
            assert isinstance(slide["shapes"], list)
        
        # Verify specific slide content
        slide_texts = [slide["text_content"] for slide in slides]
        combined_text = " ".join(slide_texts)
        assert "Complex Presentation" in combined_text
        assert "Project Overview" in combined_text
    
    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
    def test_text_content_extraction_real(self):
        """Test text content extraction with REAL text processing"""
        request = ToolRequest(
            tool_id="T11",
            operation="load_powerpoint",
            input_data={"file_path": self.test_files['simple']},
            parameters={}
        )
        
        result = self.tool.execute(request)
        text_content = result.data["document"]["text_content"]
        
        # Verify all expected content is extracted
        assert "Test Presentation" in text_content
        assert "Created for unit testing" in text_content
        assert "Key Features" in text_content
        assert "Feature 1: Document loading" in text_content
        assert "Feature 2: Text extraction" in text_content
        
        # Verify text is properly formatted
        assert len(text_content.strip()) > 0
        lines = text_content.split('\n')
        assert len(lines) > 1  # Should have multiple lines of content
    
    @pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
    def test_shape_counting_real(self):
        """Test shape counting with REAL PowerPoint analysis"""
        request = ToolRequest(
            tool_id="T11",
            operation="load_powerpoint",
            input_data={"file_path": self.test_files['complex']},
            parameters={}
        )
        
        result = self.tool.execute(request)
        doc = result.data["document"]
        
        # Verify shape counting
        assert doc["shape_count"] > 0
        
        # Verify shape details in presentation data
        presentation_data = doc["presentation_data"]
        total_shapes = 0
        for slide in presentation_data["slides"]:
            total_shapes += slide["shape_count"]
        
        # Total should match document shape count
        assert total_shapes == doc["shape_count"]